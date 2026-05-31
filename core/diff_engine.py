import os
from docx import Document
from pptx import Presentation


class DiffEngine:
    """Extract text from various file formats for diff comparison."""

    SUPPORTED_FORMATS = {'.txt', '.md', '.html', '.htm', '.docx', '.pptx'}

    @staticmethod
    def extract_text(file_path):
        """Extract text content from a file for diff comparison.
        Returns (content, format_ext, error).
        """
        if not os.path.isfile(file_path):
            return None, None, '文件不存在'

        ext = os.path.splitext(file_path)[1].lower()

        if ext not in DiffEngine.SUPPORTED_FORMATS:
            return None, ext, f'不支持的文件格式: {ext}'

        try:
            if ext == '.docx':
                content = DiffEngine._extract_docx(file_path)
            elif ext == '.pptx':
                content = DiffEngine._extract_pptx(file_path)
            else:
                content = DiffEngine._extract_text_file(file_path)
        except Exception as e:
            return None, ext, str(e)

        return content, ext, None

    @staticmethod
    def _extract_text_file(file_path):
        for encoding in ('utf-8', 'gbk', 'gb2312', 'latin-1'):
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except (UnicodeDecodeError, UnicodeError):
                continue
        return ''

    @staticmethod
    def _extract_docx(file_path):
        doc = Document(file_path)
        lines = []
        body = doc.element.body
        for child in body:
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag == 'p':
                text = ''.join(node.text or '' for node in child.iter()
                               if node.tag.endswith('}t')).strip()
                if text:
                    lines.append(text)
            elif tag == 'tbl':
                from docx.oxml.ns import qn
                for tr in child.iter(qn('w:tr')):
                    row_texts = []
                    for tc in tr.iter(qn('w:tc')):
                        cell_text = ''.join(
                            t.text or '' for t in tc.iter() if t.tag.endswith('}t')
                        ).strip()
                        row_texts.append(cell_text)
                    line = ' | '.join(row_texts)
                    if line.strip(' |'):
                        lines.append(line)
        return '\n'.join(lines)

    @staticmethod
    def _extract_pptx(file_path):
        prs = Presentation(file_path)
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f'--- Slide {i} ---')
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                        if row_text.strip(' |'):
                            lines.append(row_text)
        return '\n'.join(lines)
